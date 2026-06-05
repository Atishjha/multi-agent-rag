# ingest.py
import uuid
import os
from vector_store import get_collection


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + chunk_size]))
        i += chunk_size - overlap
    return [c for c in chunks if c.strip()]


def ingest_documents(docs: list[dict]) -> int:
    import bm25_store
    collection = get_collection()
    total = 0

    for doc in docs:
        chunks = chunk_text(doc["text"])
        if not chunks:
            continue
        ids       = [str(uuid.uuid4()) for _ in chunks]
        metadatas = [{"source": doc["source"]} for _ in chunks]

        # 1. Vector store
        collection.add(documents=chunks, ids=ids, metadatas=metadatas)

        # 2. BM25 index (incremental, no full rebuild)
        bm25_store.add_documents(chunks, metadatas, ids)

        total += len(chunks)

    return total


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_txt(path: str) -> str:
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
        return "\n".join(p.extract_text() or "" for p in PdfReader(path).pages)
    except ImportError:
        raise RuntimeError("pypdf not installed — run: pip install pypdf")


def extract_docx(path: str) -> str:
    try:
        import docx
        return "\n".join(p.text for p in docx.Document(path).paragraphs)
    except ImportError:
        raise RuntimeError("python-docx not installed — run: pip install python-docx")


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        lines = []
        for i, slide in enumerate(Presentation(path).slides, 1):
            lines.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    except ImportError:
        raise RuntimeError("python-pptx not installed — run: pip install python-pptx")


def extract_bytes(filename: str, data: bytes) -> str:
    import tempfile
    ext = os.path.splitext(filename)[1].lower()
    extractors = {
        ".txt":  extract_txt,
        ".pdf":  extract_pdf,
        ".docx": extract_docx,
        ".pptx": extract_pptx,
        ".ppt":  extract_pptx,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported file type: {ext}")
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        return extractors[ext](tmp_path)
    finally:
        os.unlink(tmp_path)


EXTRACTORS = {
    ".txt":  extract_txt,
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".ppt":  extract_pptx,
}


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    docs_folder = "docs"
    os.makedirs(docs_folder, exist_ok=True)
    docs, skipped = [], []

    for filename in os.listdir(docs_folder):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in EXTRACTORS:
            skipped.append(filename)
            continue
        filepath = os.path.join(docs_folder, filename)
        print(f"Reading: {filename} ...", end=" ")
        try:
            text = EXTRACTORS[ext](filepath)
            if not text.strip():
                print("empty, skipped.")
                skipped.append(filename)
                continue
            docs.append({"text": text, "source": filename})
            print(f"✓  ({len(text.split())} words)")
        except Exception as e:
            print(f"error: {e}")
            skipped.append(filename)

    if skipped:
        print(f"\nSkipped: {', '.join(skipped)}")
    if docs:
        total = ingest_documents(docs)
        print(f"\nIngested {total} chunks from {len(docs)} file(s).")
    else:
        print("No readable documents found.")